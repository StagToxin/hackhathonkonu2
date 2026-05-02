import ollama
import logging
import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

logger = logging.getLogger(__name__)

# Kurulu modeller: qwen:latest (4B hızlı), qwen3.6:35b-a3b (35B MoE, kaliteli)
OLLAMA_MODEL = os.getenv("OLLAMA_MODEL", "qwen3.6:35b-a3b")


def get_ai_analysis(data: dict) -> str:
    company_name = data.get("name", "Şirket")
    summary = data.get("summary", {})

    prompt = f"""Türk finans uzmanı olarak aşağıdaki şirketin finansal verilerini analiz et ve Türkçe rapor yaz.

Şirket: {company_name}
Sektör: {data.get('sector', 'Belirtilmemiş')}
Ciro: {summary.get('revenue', 'Bilinmiyor')} TL
Net Kâr: {summary.get('netProfit', 'Bilinmiyor')} TL
Borç: {summary.get('debt', 'Bilinmiyor')} TL
Özkaynak: {summary.get('equity', 'Bilinmiyor')} TL

Aşağıdaki başlıklar altında analiz yaz:
## Güçlü Yönler
## Zayıf Yönler
## Riskler
## Likidite Durumu
## Borç/Özkaynak Dengesi
## Öneriler"""

    try:
        response = ollama.chat(
            model=OLLAMA_MODEL,
            messages=[
                {"role": "system", "content": "Sen profesyonel bir Türk finansal analistsin. Kısa, net ve yapılandırılmış raporlar yazarsın."},
                {"role": "user", "content": prompt}
            ]
        )
        return response.message.content
    except Exception as e:
        logger.error(f"Ollama hatası ({OLLAMA_MODEL}): {e}")
        return f"""## Güçlü Yönler
- {company_name} için finansal veriler sisteme yüklenmiştir.
- Sektör ortalamasına göre kârlılık pozitif seyretmektedir.

## Zayıf Yönler
- Kısa vadeli borç yönetimi iyileştirme gerektirmektedir.

## Riskler
- Döviz kuru riski izlenmelidir.

## Likidite Durumu
- Mevcut likidite oranı sektör ortalaması seviyesindedir.

## Borç/Özkaynak Dengesi
- Borç/özkaynak oranı kontrol altında tutulmalıdır.

## Öneriler
1. Tahsilat vadelerini kısaltmaya odaklanın.
2. Nakit akışı projeksiyonlarını aylık güncelleyin.
3. Kredi limitlerini diversifiye edin.

*(AI analizi şu an kullanılamıyor — Ollama model: {OLLAMA_MODEL})*"""


def generate_pptx(company_name: str, analysis: str) -> str:
    Path("outputs").mkdir(exist_ok=True)

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # Kapak slaytı
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = company_name
    slide1.placeholders[1].text = "Finansal Analiz ve Risk Raporu"

    tf = slide1.shapes.title.text_frame
    tf.paragraphs[0].runs[0].font.color.rgb = RGBColor(0x1E, 0x40, 0xAF)
    tf.paragraphs[0].runs[0].font.size = Pt(36)

    # Analiz slaytları — her başlık ayrı slide
    sections = []
    current_title = ""
    current_body = []

    for line in analysis.split("\n"):
        if line.startswith("## "):
            if current_title:
                sections.append((current_title, "\n".join(current_body).strip()))
            current_title = line[3:].strip()
            current_body = []
        else:
            current_body.append(line)

    if current_title:
        sections.append((current_title, "\n".join(current_body).strip()))

    for title, body in sections:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = title
        slide.shapes.title.text_frame.paragraphs[0].runs[0].font.size = Pt(24)
        slide.shapes.title.text_frame.paragraphs[0].runs[0].font.bold = True

        tf = slide.placeholders[1].text_frame
        tf.text = body if body else "-"
        for para in tf.paragraphs:
            for run in para.runs:
                run.font.size = Pt(14)

    safe_name = "".join(c if c.isalnum() or c in " _-" else "_" for c in company_name)
    output_path = f"outputs/{safe_name}_sunum.pptx"
    prs.save(output_path)
    return output_path


def parse_ocr_with_ai(text: str) -> dict:
    prompt = f"""Aşağıdaki belgeden şirket bilgilerini JSON formatında çıkar.
Sadece JSON döndür, başka hiçbir şey yazma.

Belge içeriği:
{text[:2000]}

Çıkarılacak alanlar:
{{
  "name": "şirket adı",
  "taxNo": "vergi numarası",
  "tradeRegistryNo": "ticaret sicil numarası",
  "sector": "sektör",
  "contactName": "yetkili kişi adı",
  "phone": "telefon",
  "email": "e-posta",
  "address": "adres",
  "foundedAt": "kuruluş tarihi (YYYY-MM-DD)",
  "estimatedTurnover": 0,
  "contractType": "sözleşme türü"
}}"""

    try:
        response = ollama.chat(
            model=OLLAMA_MODEL,
            messages=[{"role": "user", "content": prompt}]
        )
        content = response.message.content.strip()
        # JSON bloğunu bul
        import json, re
        match = re.search(r'\{.*\}', content, re.DOTALL)
        if match:
            return json.loads(match.group())
    except Exception as e:
        logger.error(f"OCR AI hatası: {e}")

    return {
        "name": "", "taxNo": "", "tradeRegistryNo": "", "sector": "",
        "contactName": "", "phone": "", "email": "", "address": "",
        "foundedAt": "", "estimatedTurnover": 0, "contractType": "Analiz"
    }
